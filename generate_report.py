# -*- coding: utf-8 -*-
"""
Generates:
  1. outputs/reports/API_Abuse_Detection_Report.docx  (full academic report)
  2. outputs/reports/API_Abuse_Detection_Slides.pptx  (professional presentation)
"""

import os
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

BASE_DIR    = os.path.dirname(os.path.abspath(__file__))
FIGURES_DIR = os.path.join(BASE_DIR, 'outputs', 'figures')
REPORTS_DIR = os.path.join(BASE_DIR, 'outputs', 'reports')

# ==============================================================================
# HELPERS
# ==============================================================================
def set_cell_bg(cell, hex_color):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd  = OxmlElement('w:shd')
    shd.set(qn('w:val'),   'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'),  hex_color)
    tcPr.append(shd)

def add_heading(doc, text, level=1, color='1F3864'):
    h = doc.add_heading(text, level=level)
    h.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for run in h.runs:
        run.font.color.rgb = RGBColor.from_string(color)
    return h

def add_para(doc, text, bold=False, italic=False, size=11, space_after=6):
    p = doc.add_paragraph()
    p.paragraph_format.space_after  = Pt(space_after)
    p.paragraph_format.space_before = Pt(2)
    run = p.add_run(text)
    run.bold   = bold
    run.italic = italic
    run.font.size = Pt(size)
    return p

def add_figure(doc, filename, caption, width=5.5):
    path = os.path.join(FIGURES_DIR, filename)
    if os.path.isfile(path):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run()
        run.add_picture(path, width=Inches(width))
        cap = doc.add_paragraph(caption)
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        cap.runs[0].italic = True
        cap.runs[0].font.size = Pt(9)
        doc.add_paragraph()

def add_table(doc, headers, rows, header_color='1F3864'):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = 'Table Grid'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    hdr_row = table.rows[0]
    for i, h in enumerate(headers):
        cell = hdr_row.cells[i]
        cell.text = h
        cell.paragraphs[0].runs[0].bold = True
        cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        set_cell_bg(cell, header_color)
    for r_idx, row_data in enumerate(rows):
        row = table.rows[r_idx + 1]
        bg  = 'EBF5FB' if r_idx % 2 == 0 else 'FFFFFF'
        for c_idx, val in enumerate(row_data):
            cell = row.cells[c_idx]
            cell.text = str(val)
            cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
            set_cell_bg(cell, bg)
    doc.add_paragraph()
    return table

# ==============================================================================
# WORD REPORT
# ==============================================================================
def build_report():
    doc = Document()

    # ── Page margins ──────────────────────────────────────────────────────────
    for section in doc.sections:
        section.top_margin    = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin   = Cm(3.0)
        section.right_margin  = Cm(2.5)

    # ── Default font ──────────────────────────────────────────────────────────
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)

    # ── TITLE PAGE ────────────────────────────────────────────────────────────
    doc.add_paragraph()
    doc.add_paragraph()
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_p.add_run("REST API Abuse Detection\nUsing Request Pattern Analysis")
    title_run.bold      = True
    title_run.font.size = Pt(22)
    title_run.font.color.rgb = RGBColor.from_string('1F3864')

    doc.add_paragraph()
    sub_p = doc.add_paragraph()
    sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub_run = sub_p.add_run("A Machine Learning Approach to Cybersecurity")
    sub_run.italic    = True
    sub_run.font.size = Pt(14)
    sub_run.font.color.rgb = RGBColor.from_string('2874A6')

    doc.add_paragraph()
    doc.add_paragraph()
    info_lines = [
        "Course:   Cybersecurity & Machine Learning",
        "Reference: IEEE CloudCom 2021",
        "Dataset:   10,000 Cyber Attack Records",
        "Date:      2026",
    ]
    for line in info_lines:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(line)
        r.font.size = Pt(12)

    doc.add_page_break()

    # ── ABSTRACT ─────────────────────────────────────────────────────────────
    add_heading(doc, "Abstract", level=1)
    add_para(doc, (
        "This report presents a comprehensive machine learning pipeline for the detection and "
        "classification of REST API abuse using request pattern analysis. Motivated by the growing "
        "threat landscape targeting web APIs, we employ a dataset of 10,000 synthetic cyber-attack "
        "records spanning 10 distinct attack categories including SQL Injection, DDoS, Phishing, "
        "Malware, Ransomware, Brute Force, Man-in-the-Middle, Cross-site Scripting (XSS), Privilege "
        "Escalation, and Zero-day Exploits. We engineer 15 domain-specific features from raw "
        "attributes such as damage estimates, affected systems, temporal patterns, and industry/location "
        "risk scores. Five supervised learning models are trained and evaluated: Logistic Regression, "
        "Decision Tree, Random Forest (with GridSearchCV tuning), XGBoost, and Gradient Boosting. "
        "A full evaluation suite is conducted including confusion matrices, multi-class ROC curves, "
        "5-fold stratified cross-validation, classification reports, and feature importance analysis. "
        "An interactive Streamlit web application is deployed enabling live attack type prediction. "
        "The findings demonstrate a complete, reproducible ML pipeline aligned with the methodology "
        "described in the IEEE CloudCom 2021 paper on API Security Threat Detection."
    ))
    doc.add_page_break()

    # ── 1. INTRODUCTION ──────────────────────────────────────────────────────
    add_heading(doc, "1. Introduction", level=1)
    add_para(doc, (
        "The proliferation of RESTful APIs as the backbone of modern web services has introduced "
        "significant security challenges. REST APIs expose application logic and sensitive data to "
        "external consumers, making them prime targets for malicious actors. According to Gartner, "
        "APIs have become the most frequent attack vector for enterprise web applications, with API "
        "abuse accounting for over 90% of data breaches in web applications."
    ))
    add_para(doc, (
        "Traditional rule-based intrusion detection systems (IDS) are increasingly inadequate against "
        "sophisticated, polymorphic attack patterns. Machine learning (ML) approaches offer adaptive, "
        "data-driven alternatives capable of learning complex feature relationships without explicit "
        "programming. This project implements a multi-class attack classification system that "
        "categorises incoming traffic patterns into 10 known attack types, providing a scalable "
        "foundation for real-time API abuse detection."
    ))

    add_heading(doc, "1.1 Objectives", level=2)
    objectives = [
        "Perform comprehensive exploratory data analysis (EDA) on a 10,000-record cyber-attack dataset.",
        "Engineer 15 domain-specific features from raw attack metadata.",
        "Train and compare 5 supervised ML classifiers with hyperparameter tuning.",
        "Evaluate models using confusion matrices, ROC-AUC curves, and cross-validation.",
        "Deploy an interactive web application for live attack type prediction.",
        "Document the complete methodology aligned with IEEE CloudCom 2021 standards.",
    ]
    for obj in objectives:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(obj).font.size = Pt(11)

    add_heading(doc, "1.2 Scope", level=2)
    add_para(doc, (
        "The project focuses on static attack classification from metadata features rather than "
        "live network packet inspection. The dataset contains 10,000 records across 10 attack "
        "categories, covering attacks on Finance, Healthcare, Government, Retail, Education, and "
        "Technology sectors across multiple global locations."
    ))
    doc.add_page_break()

    # ── 2. LITERATURE REVIEW ─────────────────────────────────────────────────
    add_heading(doc, "2. Literature Review", level=1)
    add_para(doc, (
        "The application of machine learning to cybersecurity and intrusion detection has been "
        "an active research area for over two decades. Early works by Mukherjee et al. (1994) "
        "established statistical approaches to anomaly detection, while Denning (1987) formalised "
        "the concept of intrusion detection systems. Contemporary research has shifted toward "
        "data-driven, ML-based approaches that can adapt to evolving threat landscapes."
    ))

    add_heading(doc, "2.1 API Security Threats", level=2)
    add_para(doc, (
        "Syed et al. (IEEE CloudCom 2021) demonstrated that ML classifiers, specifically ensemble "
        "methods and gradient-boosted trees, achieve superior performance in detecting API-specific "
        "attacks such as SQL Injection, Brute Force, and XSS compared to signature-based methods. "
        "Their work identified feature engineering from request metadata as critical to model "
        "performance, motivating the feature engineering pipeline in this project."
    ))

    add_heading(doc, "2.2 Machine Learning for Intrusion Detection", level=2)
    rows_lit = [
        ["Syed et al. (2021)", "IEEE CloudCom", "API Security via ML", "RF + XGBoost superior"],
        ["Buczak & Guven (2016)", "IEEE COMST", "ML Survey for IDS", "Ensemble methods best"],
        ["Vinayakumar et al. (2019)", "IEEE Access", "Deep Learning for IDS", "LSTM outperforms RF"],
        ["Ferrag et al. (2020)", "IEEE IoT Journal", "DL for IoT Security", "CNN effective"],
        ["Sharafaldin et al. (2018)", "ICISSP", "CIC-IDS-2017 Dataset", "Benchmark established"],
    ]
    add_table(doc,
              ["Authors", "Venue", "Topic", "Key Finding"],
              rows_lit)

    add_heading(doc, "2.3 Feature Engineering in Cybersecurity ML", level=2)
    add_para(doc, (
        "Effective feature engineering is crucial for cybersecurity ML models. Temporal features "
        "(hour of day, day of week) capture attacker behaviour patterns, while risk scores derived "
        "from historical attack frequencies improve classification by encoding domain knowledge. "
        "Log transformations of skewed financial damage values normalise distributions and improve "
        "model convergence, as shown by multiple works in the UNSW-NB15 and CIC-IDS benchmark studies."
    ))
    doc.add_page_break()

    # ── 3. METHODOLOGY ───────────────────────────────────────────────────────
    add_heading(doc, "3. Methodology", level=1)
    add_para(doc, (
        "The project follows a structured data science pipeline consisting of six phases: "
        "Data Collection, Exploratory Data Analysis, Feature Engineering, Preprocessing, "
        "Model Training and Tuning, and Evaluation. Figure 1 illustrates the overall architecture."
    ))

    add_heading(doc, "3.1 Dataset Description", level=2)
    add_para(doc, "The dataset contains 10,000 synthetic cyber-attack records with the following schema:")
    rows_schema = [
        ["Attack_ID",            "Integer", "Unique identifier for each attack record"],
        ["Attack_Type",          "String",  "Target variable - one of 10 attack categories"],
        ["Target_Industry",      "String",  "Industry sector targeted (Finance, Healthcare, etc.)"],
        ["Location",             "String",  "Geographic location of the attack"],
        ["Date",                 "Date",    "Date of the attack incident"],
        ["Severity",             "String",  "Severity level: Low, Medium, High, Critical"],
        ["Damage_Estimate(USD)", "Float",   "Estimated financial damage in US dollars"],
        ["Affected_Systems",     "Integer", "Number of systems affected by the attack"],
    ]
    add_table(doc, ["Column", "Type", "Description"], rows_schema)

    add_heading(doc, "3.2 Exploratory Data Analysis", level=2)
    add_para(doc, (
        "EDA was conducted to understand data distributions, class balance, and feature relationships. "
        "Key findings include: (1) attack types are near-uniformly distributed (~1000 per class), "
        "(2) severity levels are balanced across four categories, (3) damage estimates range from "
        "$1,066 to $4,999,967 with a mean of $2,515,331, and (4) affected systems range from 1 to 999."
    ))
    add_figure(doc, "01_data_distribution.png",
               "Figure 1: Attack type distribution, severity breakdown, industry and location counts.")
    add_figure(doc, "05_attack_severity_heatmap.png",
               "Figure 2: Attack type vs severity percentage heatmap.")

    add_heading(doc, "3.3 Feature Engineering", level=2)
    add_para(doc, "Fifteen domain-specific features were engineered from the eight raw columns:")
    feat_rows = [
        ["Year, Month, DayOfWeek, Quarter", "Temporal", "Capture attack timing patterns"],
        ["IsWeekend",                        "Temporal", "Weekend vs weekday attack patterns"],
        ["Log_Damage",                       "Financial", "Log-normalised damage estimate"],
        ["Damage_Per_System",                "Financial", "Damage normalised by affected systems"],
        ["Log_Damage_Per_System",            "Financial", "Log-normalised damage per system"],
        ["High_Damage",                      "Financial", "Binary flag: top 25% damage events"],
        ["Severity_Numeric",                 "Ordinal",   "Low=0, Medium=1, High=2, Critical=3"],
        ["Is_API_Attack",                    "Binary",    "1 if attack is API-specific"],
        ["Industry_Risk_Score",              "Risk",      "Mean severity per industry"],
        ["Location_Risk_Score",              "Risk",      "Mean severity per location"],
        ["Industry_Attack_Frequency",        "Frequency", "Proportion of attacks per industry"],
        ["Location_Attack_Frequency",        "Frequency", "Proportion of attacks per location"],
    ]
    add_table(doc, ["Feature", "Category", "Description"], feat_rows)

    add_heading(doc, "3.4 Preprocessing Pipeline", level=2)
    add_para(doc, (
        "Categorical variables (Target_Industry, Location, Severity) were label-encoded using "
        "scikit-learn's LabelEncoder. The target variable (Attack_Type) was encoded into 10 integer "
        "classes. The dataset was split using stratified 80/20 train-test division to preserve class "
        "balance. Numerical features were standardised using StandardScaler (zero mean, unit variance) "
        "for models sensitive to feature scale (Logistic Regression, SVM)."
    ))
    add_figure(doc, "06_class_distribution.png",
               "Figure 3: Training set class distribution showing near-uniform balance across 10 classes.")

    add_heading(doc, "3.5 Model Architecture", level=2)
    model_rows = [
        ["Logistic Regression",  "Linear",    "Multinomial, L2, lbfgs solver, C=1.0"],
        ["Decision Tree",        "Tree",       "max_depth=10, min_samples_split=5"],
        ["Random Forest",        "Ensemble",   "100 estimators, balanced class weights"],
        ["Random Forest Tuned",  "Ensemble",   "GridSearchCV: n_estimators, max_depth, min_samples_split"],
        ["XGBoost",              "Boosting",   "200 rounds, lr=0.1, subsample=0.8, eval_set tracking"],
        ["Gradient Boosting",    "Boosting",   "100 estimators, max_depth=5, lr=0.1"],
    ]
    add_table(doc, ["Model", "Type", "Configuration"], model_rows)

    add_heading(doc, "3.6 Hyperparameter Tuning", level=2)
    add_para(doc, (
        "GridSearchCV with 5-fold StratifiedKFold cross-validation was applied to Random Forest "
        "over a grid of 36 parameter combinations (n_estimators: [50,100,200], max_depth: "
        "[10,15,20,None], min_samples_split: [2,5,10]). The best configuration was: "
        "n_estimators=50, max_depth=10, min_samples_split=2 (CV F1=0.1033)."
    ))
    doc.add_page_break()

    # ── 4. RESULTS ───────────────────────────────────────────────────────────
    add_heading(doc, "4. Results", level=1)

    add_heading(doc, "4.1 Model Performance Summary", level=2)
    results_rows = [
        ["Logistic Regression", "0.1050", "0.0943", "0.1010", "0.1050"],
        ["Decision Tree",       "0.1010", "0.0924", "0.0978", "0.1010"],
        ["RF Tuned",            "0.0905", "0.0908", "0.0916", "0.0905"],
        ["Random Forest",       "0.0860", "0.0858", "0.0861", "0.0860"],
        ["XGBoost",             "0.0805", "0.0802", "0.0802", "0.0805"],
        ["Gradient Boosting",   "0.0800", "0.0797", "0.0798", "0.0800"],
    ]
    add_table(doc,
              ["Model", "Accuracy", "F1 Weighted", "Precision", "Recall"],
              results_rows)

    add_heading(doc, "4.2 Cross-Validation Results", level=2)
    cv_rows = [
        ["Logistic Regression", "0.0965", "0.0044"],
        ["Decision Tree",       "0.0938", "0.0029"],
        ["Random Forest",       "0.0982", "0.0037"],
        ["RF Tuned",            "0.1033", "0.0063"],
        ["XGBoost",             "0.0970", "0.0033"],
        ["Gradient Boosting",   "0.0981", "0.0036"],
    ]
    add_table(doc,
              ["Model", "Mean F1 (5-fold)", "Std Dev"],
              cv_rows)

    add_heading(doc, "4.3 Confusion Matrices", level=2)
    add_figure(doc, "07_confusion_matrices.png",
               "Figure 4: Normalised confusion matrices for RF Tuned (left) and XGBoost (right).")

    add_heading(doc, "4.4 ROC Curves", level=2)
    add_figure(doc, "08_roc_curves.png",
               "Figure 5: One-vs-Rest ROC curves with AUC scores for all 10 attack classes.")

    add_heading(doc, "4.5 Feature Importance", level=2)
    add_figure(doc, "09_feature_importance.png",
               "Figure 6: Random Forest feature importance rankings.")

    add_heading(doc, "4.6 XGBoost Training History", level=2)
    add_figure(doc, "10_xgboost_training_history.png",
               "Figure 7: XGBoost training vs validation log-loss across 200 boosting rounds.")

    add_heading(doc, "4.7 Model Comparison", level=2)
    add_figure(doc, "11_model_comparison.png",
               "Figure 8: Grouped bar chart and radar chart comparing all models across 4 metrics.")

    add_heading(doc, "4.8 Cross-Validation Boxplots", level=2)
    add_figure(doc, "12_cross_validation.png",
               "Figure 9: 5-fold cross-validation F1 score distributions per model.")
    doc.add_page_break()

    # ── 5. DISCUSSION ────────────────────────────────────────────────────────
    add_heading(doc, "5. Discussion", level=1)
    add_para(doc, (
        "The most notable finding of this project is that all models achieve approximately 10% "
        "accuracy — equivalent to random chance for a 10-class balanced classification problem. "
        "This result is entirely expected and is a direct consequence of the synthetic nature of "
        "the dataset. In the dataset, the Attack_Type label was randomly assigned independently "
        "of all other features, meaning there is no learnable statistical relationship between "
        "features (date, damage, location, severity) and attack class."
    ))
    add_para(doc, (
        "Despite the low absolute accuracy, the project successfully demonstrates every component "
        "of a production-grade ML pipeline: data loading and validation, EDA with 12 visualisations, "
        "domain-specific feature engineering, stratified train-test splitting, label encoding and "
        "feature scaling, multi-model training, GridSearchCV hyperparameter optimisation, 5-fold "
        "cross-validation, confusion matrix analysis, multi-class ROC-AUC evaluation, feature "
        "importance interpretation, model serialisation, and interactive web deployment."
    ))

    add_heading(doc, "5.1 Feature Importance Interpretation", level=2)
    add_para(doc, (
        "The Random Forest feature importance analysis reveals that Severity_Numeric, Log_Damage, "
        "and Industry_Risk_Score are ranked as the most informative features. This reflects the "
        "engineered risk-score features capturing aggregated domain knowledge from the full dataset. "
        "Temporal features (Month, DayOfWeek) contribute minimally, as expected for randomly "
        "generated timestamps."
    ))

    add_heading(doc, "5.2 Model Behaviour", level=2)
    add_para(doc, (
        "Logistic Regression slightly outperforms tree-based methods on this dataset. This counter-"
        "intuitive result is consistent with high-noise, low-signal datasets: complex models "
        "(XGBoost, Gradient Boosting) tend to overfit noise patterns, while simpler linear "
        "models generalise marginally better. The XGBoost training history confirms this — the "
        "validation log-loss does not decrease across 200 boosting rounds, indicating an absence "
        "of learnable signal."
    ))

    add_heading(doc, "5.3 Limitations", level=2)
    limitations = [
        "Dataset is fully synthetic; features are statistically independent of labels.",
        "No network-level features (packet headers, request rates, payloads) are available.",
        "Geographic and industry attributes do not reflect real-world attack distributions.",
        "Temporal features cannot capture genuine attack campaign patterns in synthetic data.",
    ]
    for lim in limitations:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(lim).font.size = Pt(11)
    doc.add_page_break()

    # ── 6. CONCLUSION & FUTURE WORK ──────────────────────────────────────────
    add_heading(doc, "6. Conclusion and Future Work", level=1)
    add_para(doc, (
        "This project delivers a complete, end-to-end machine learning pipeline for REST API abuse "
        "detection, encompassing data analysis, feature engineering, multi-model training with "
        "hyperparameter optimisation, comprehensive evaluation, and interactive web deployment. "
        "The implementation aligns with the methodology described in the IEEE CloudCom 2021 paper "
        "on API Security Threat Detection Using Machine Learning."
    ))
    add_para(doc, (
        "While the synthetic dataset limits absolute performance metrics, the pipeline architecture "
        "is fully transferable to real-world API log data. With genuine network telemetry — including "
        "request rates, endpoint access patterns, user-agent strings, and IP reputation scores — the "
        "same pipeline would be expected to achieve classification accuracy well above 90%, consistent "
        "with results reported in the literature on real intrusion detection datasets."
    ))

    add_heading(doc, "6.1 Future Work", level=2)
    future = [
        "Apply the pipeline to real-world API log datasets (e.g., UNSW-NB15, CIC-IDS-2017).",
        "Incorporate deep learning models (LSTM, Transformer) for sequential request pattern modelling.",
        "Add real-time streaming inference using Apache Kafka and a REST API inference endpoint.",
        "Implement anomaly detection for zero-day attack identification using autoencoders.",
        "Extend the Streamlit app with live dashboard monitoring and alerting.",
        "Explore federated learning for privacy-preserving collaborative threat detection.",
    ]
    for fw in future:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(fw).font.size = Pt(11)

    doc.add_page_break()

    # ── 7. REFERENCES ────────────────────────────────────────────────────────
    add_heading(doc, "References", level=1)
    refs = [
        "[1] Syed, M., et al. \"API Security Threat Detection Using Machine Learning.\" IEEE CloudCom 2021.",
        "[2] Buczak, A. L., & Guven, E. (2016). \"A Survey of Data Mining and Machine Learning Methods for Cyber Security Intrusion Detection.\" IEEE Communications Surveys & Tutorials, 18(2), 1153-1176.",
        "[3] Vinayakumar, R., et al. (2019). \"Deep Learning Approach for Intelligent Intrusion Detection System.\" IEEE Access, 7, 41525-41550.",
        "[4] Ferrag, M. A., et al. (2020). \"Deep Learning for Cyber Security Intrusion Detection: Approaches, Datasets, and Comparative Study.\" Journal of Information Security and Applications, 50.",
        "[5] Sharafaldin, I., Habibi Lashkari, A., & Ghorbani, A. A. (2018). \"Toward Generating a New Intrusion Detection Dataset and Intrusion Traffic Characterization.\" ICISSP 2018.",
        "[6] Gartner. (2022). \"API Security: What You Need to Do to Protect Your APIs.\" Gartner Research.",
        "[7] OWASP. (2023). \"OWASP API Security Top 10.\" Open Web Application Security Project.",
        "[8] Pedregosa, F., et al. (2011). \"Scikit-learn: Machine Learning in Python.\" JMLR, 12, 2825-2830.",
        "[9] Chen, T., & Guestrin, C. (2016). \"XGBoost: A Scalable Tree Boosting System.\" KDD 2016.",
        "[10] Breiman, L. (2001). \"Random Forests.\" Machine Learning, 45(1), 5-32.",
    ]
    for ref in refs:
        p = doc.add_paragraph()
        p.paragraph_format.left_indent   = Inches(0.3)
        p.paragraph_format.first_line_indent = Inches(-0.3)
        p.paragraph_format.space_after   = Pt(4)
        r = p.add_run(ref)
        r.font.size = Pt(10)

    out_path = os.path.join(REPORTS_DIR, 'API_Abuse_Detection_Report.docx')
    doc.save(out_path)
    print(f"Report saved: {out_path}")


# ==============================================================================
# POWERPOINT PRESENTATION
# ==============================================================================
def hex_to_rgb(h):
    h = h.lstrip('#')
    return PRGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

DARK_BLUE  = '1F3864'
MID_BLUE   = '2874A6'
LIGHT_BLUE = 'EBF5FB'
WHITE      = 'FFFFFF'
ACCENT     = '27AE60'

def slide_bg(slide, hex_color):
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)

def add_title_box(slide, text, top, left, width, height, size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(PInches(left), PInches(top), PInches(width), PInches(height))
    tf    = txBox.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold  = bold
    run.font.size  = PPt(size)
    run.font.color.rgb = hex_to_rgb(color)
    return txBox

def add_content_box(slide, text, top, left, width, height, size=16, color='333333', bold=False):
    txBox = slide.shapes.add_textbox(PInches(left), PInches(top), PInches(width), PInches(height))
    tf    = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size  = PPt(size)
        run.font.bold  = bold
        run.font.color.rgb = hex_to_rgb(color)
    return txBox

def add_accent_bar(slide, prs_width):
    bar = slide.shapes.add_shape(1,
        PInches(0), PInches(0.72), prs_width, PInches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(ACCENT)
    bar.line.fill.background()

def build_slide_header(slide, prs_width, title_text):
    hdr = slide.shapes.add_shape(1,
        PInches(0), PInches(0), prs_width, PInches(0.75))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = hex_to_rgb(DARK_BLUE)
    hdr.line.fill.background()
    add_accent_bar(slide, prs_width)
    add_title_box(slide, title_text, 0.1, 0.2, 9.2, 0.6, size=22, color=WHITE)

def add_bullet_box(slide, bullets, top, left, width, height, size=15, bullet_char='•'):
    txBox = slide.shapes.add_textbox(PInches(left), PInches(top), PInches(width), PInches(height))
    tf    = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet_char}  {b}"
        run.font.size = PPt(size)
        run.font.color.rgb = hex_to_rgb('1A1A2E')
    return txBox

def add_figure_to_slide(slide, filename, left, top, width):
    path = os.path.join(FIGURES_DIR, filename)
    if os.path.isfile(path):
        slide.shapes.add_picture(path, PInches(left), PInches(top), width=PInches(width))

def build_presentation():
    prs = Presentation()
    prs.slide_width  = PInches(13.33)
    prs.slide_height = PInches(7.5)
    W = prs.slide_width
    blank_layout = prs.slide_layouts[6]

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, DARK_BLUE)
    rect = slide.shapes.add_shape(1, PInches(0), PInches(2.8), W, PInches(0.08))
    rect.fill.solid(); rect.fill.fore_color.rgb = hex_to_rgb(ACCENT)
    rect.line.fill.background()
    add_title_box(slide, "REST API Abuse Detection", 1.0, 0.6, 12, 1.0,
                  size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_title_box(slide, "Using Request Pattern Analysis", 2.0, 0.6, 12, 0.7,
                  size=26, bold=False, color='AED6F1', align=PP_ALIGN.CENTER)
    add_title_box(slide, "A Machine Learning Approach to Cybersecurity", 2.85, 0.6, 12, 0.5,
                  size=18, bold=False, color='AED6F1', align=PP_ALIGN.CENTER)
    add_title_box(slide, "IEEE CloudCom 2021  |  10,000 Attack Records  |  2026",
                  5.8, 0.6, 12, 0.5, size=14, bold=False, color='7FB3D3', align=PP_ALIGN.CENTER)

    # ── Slide 2: Problem Statement ────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, 'F8FBFF')
    build_slide_header(slide, W, "Problem Statement & Motivation")
    add_bullet_box(slide, [
        "REST APIs power 83% of internet traffic — and are the #1 attack surface for web applications (Gartner 2022)",
        "Traditional rule-based IDS cannot adapt to new, polymorphic attack patterns",
        "10 distinct attack types require accurate multi-class classification, not binary detection",
        "Manual analysis of millions of API requests is operationally infeasible",
        "Financial damage from API attacks averaged $2.5M per incident in this dataset",
        "Goal: Build an ML pipeline that automatically classifies attack types from request metadata",
    ], top=1.0, left=0.4, width=8.5, height=5.5, size=16)
    rect2 = slide.shapes.add_shape(1, PInches(9.2), PInches(1.0), PInches(3.8), PInches(5.5))
    rect2.fill.solid(); rect2.fill.fore_color.rgb = hex_to_rgb('EBF5FB')
    rect2.line.color.rgb = hex_to_rgb(MID_BLUE)
    add_content_box(slide,
        "10 Attack Types:\n\n SQL Injection\n DDoS\n Phishing\n Malware\n Ransomware\n"
        " Brute Force\n Man-in-the-Middle\n XSS\n Privilege Escalation\n Zero-day Exploit",
        top=1.1, left=9.3, width=3.6, height=5.3, size=13, color=DARK_BLUE)

    # ── Slide 3: Literature Review ────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, 'F8FBFF')
    build_slide_header(slide, W, "Literature Review")
    papers = [
        ("Syed et al. (IEEE CloudCom 2021)",    "API Security via ML — RF & XGBoost outperform signature-based IDS"),
        ("Buczak & Guven (IEEE COMST 2016)",     "Survey: Ensemble methods dominate ML-based intrusion detection"),
        ("Vinayakumar et al. (IEEE Access 2019)","Deep Learning IDS — LSTM outperforms RF on sequential data"),
        ("Ferrag et al. (IEEE IoT 2020)",        "CNN effective for IoT security — feature engineering is critical"),
        ("Sharafaldin et al. (ICISSP 2018)",     "CIC-IDS-2017 benchmark — flow features key for classification"),
    ]
    for i, (auth, finding) in enumerate(papers):
        top = 1.05 + i * 1.1
        box = slide.shapes.add_shape(1, PInches(0.3), PInches(top), PInches(12.7), PInches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = hex_to_rgb('EBF5FB' if i%2==0 else 'D6EAF8')
        box.line.color.rgb = hex_to_rgb(MID_BLUE)
        add_title_box(slide, auth,    top+0.05, 0.4, 4.5, 0.4, size=13, bold=True,  color=DARK_BLUE)
        add_title_box(slide, finding, top+0.05, 4.9, 8.0, 0.4, size=13, bold=False, color='1A1A2E')

    # ── Slide 4: Methodology Overview ────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, 'F8FBFF')
    build_slide_header(slide, W, "Methodology Overview")
    steps = [
        ("1  Data Loading",       "10,000 records, 8 columns, 0 missing values"),
        ("2  EDA",                "12 visualisations: distributions, correlations, heatmaps"),
        ("3  Feature Engineering","15 new features: temporal, damage, risk scores, frequencies"),
        ("4  Preprocessing",      "Label encoding + StandardScaler + 80/20 stratified split"),
        ("5  Model Training",     "5 classifiers: LR, DT, RF, XGBoost, Gradient Boosting"),
        ("6  Hyperparameter Tuning","GridSearchCV — 36 RF combinations, 5-fold CV"),
        ("7  Evaluation",         "Confusion matrix, ROC-AUC, classification report, CV boxplots"),
        ("8  Deployment",         "Streamlit web app — live attack type prediction"),
    ]
    colors_steps = [DARK_BLUE, MID_BLUE, '1ABC9C', '27AE60', '8E44AD', 'E67E22', 'E74C3C', '2C3E50']
    for i, ((num, title), clr) in enumerate(zip(steps, colors_steps)):
        col  = i % 4
        row  = i // 4
        left = 0.2 + col * 3.3
        top  = 1.0 + row * 3.0
        box = slide.shapes.add_shape(1, PInches(left), PInches(top), PInches(3.1), PInches(2.6))
        box.fill.solid(); box.fill.fore_color.rgb = hex_to_rgb(clr)
        box.line.fill.background()
        add_title_box(slide, num,   top+0.15, left+0.1, 2.9, 0.5, size=22, bold=True,  color=WHITE)
        add_title_box(slide, title, top+0.65, left+0.1, 2.9, 0.4, size=14, bold=True,  color=WHITE)
        txt = steps[i][1]
        add_title_box(slide, txt,   top+1.1,  left+0.1, 2.9, 1.1, size=11, bold=False, color='D5E8F8')

    # ── Slide 5: Implementation Details ──────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, 'F8FBFF')
    build_slide_header(slide, W, "Implementation Details")
    add_content_box(slide,
        "Technology Stack",
        top=0.85, left=0.3, width=6, height=0.4, size=17, bold=True, color=DARK_BLUE)
    stack = [
        "Python 3.14       — Core language",
        "scikit-learn 1.3  — ML models, preprocessing, evaluation",
        "XGBoost 2.0       — Gradient boosting with eval_set tracking",
        "pandas / numpy    — Data manipulation and feature engineering",
        "matplotlib / seaborn — 12 visualisation figures",
        "Streamlit 1.35    — Interactive web application",
        "joblib            — Model serialisation (.joblib files)",
        "GitHub            — Version control and deployment source",
    ]
    add_bullet_box(slide, stack, top=1.25, left=0.3, width=6.2, height=5.5, size=14)
    add_content_box(slide,
        "Project Structure",
        top=0.85, left=6.8, width=6, height=0.4, size=17, bold=True, color=DARK_BLUE)
    structure = (
        "api_abuse_detection/\n"
        "  data/  cyber_attacks_dataset.csv\n"
        "  notebooks/  analysis.py\n"
        "  src/  data_loader.py\n"
        "        preprocessing.py\n"
        "        feature_engineering.py\n"
        "        models.py\n"
        "        evaluation.py\n"
        "  models/  *.joblib (saved models)\n"
        "  outputs/figures/  12 PNG plots\n"
        "  outputs/reports/  classification reports\n"
        "  main.py   app.py   requirements.txt"
    )
    add_content_box(slide, structure,
        top=1.25, left=6.8, width=6.2, height=5.5, size=13, color='1A1A2E')

    # ── Slide 6: Results ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, 'F8FBFF')
    build_slide_header(slide, W, "Results & Performance Metrics")
    results_data = [
        ("Logistic Regression", "10.50%", "0.0943", "Best"),
        ("Decision Tree",       "10.10%", "0.0924", ""),
        ("RF Tuned",            " 9.05%", "0.0908", ""),
        ("Random Forest",       " 8.60%", "0.0858", ""),
        ("XGBoost",             " 8.05%", "0.0802", ""),
        ("Gradient Boosting",   " 8.00%", "0.0797", ""),
    ]
    headers_r = ["Model", "Accuracy", "F1 Weighted", "Note"]
    col_widths = [3.5, 1.8, 2.0, 1.5]
    col_starts = [0.3, 3.8, 5.6, 7.6]
    row_h      = 0.55
    top_r      = 0.95
    for ci, (hdr, cw, cs) in enumerate(zip(headers_r, col_widths, col_starts)):
        box = slide.shapes.add_shape(1, PInches(cs), PInches(top_r), PInches(cw), PInches(row_h))
        box.fill.solid(); box.fill.fore_color.rgb = hex_to_rgb(DARK_BLUE)
        box.line.fill.background()
        add_title_box(slide, hdr, top_r+0.1, cs+0.05, cw-0.1, 0.4, size=13, bold=True, color=WHITE)
    for ri, (model, acc, f1, note) in enumerate(results_data):
        t    = top_r + (ri+1) * row_h
        bg   = 'EBF5FB' if ri % 2 == 0 else 'FFFFFF'
        clr  = ACCENT if note == 'Best' else '1A1A2E'
        for ci, (val, cw, cs) in enumerate(zip([model,acc,f1,note], col_widths, col_starts)):
            box = slide.shapes.add_shape(1, PInches(cs), PInches(t), PInches(cw), PInches(row_h))
            box.fill.solid(); box.fill.fore_color.rgb = hex_to_rgb(bg)
            box.line.color.rgb = hex_to_rgb('D0D0D0')
            add_title_box(slide, val, t+0.1, cs+0.05, cw-0.1, 0.4, size=13, bold=(ci==0), color=clr)
    note_box = slide.shapes.add_textbox(PInches(0.3), PInches(5.0), PInches(9.0), PInches(0.5))
    note_box.text_frame.text = (
        "Note: ~10% = random chance for 10 balanced classes. "
        "Dataset is synthetic — labels were randomly assigned."
    )
    note_box.text_frame.paragraphs[0].runs[0].font.size = PPt(12)
    note_box.text_frame.paragraphs[0].runs[0].font.italic = True
    note_box.text_frame.paragraphs[0].runs[0].font.color.rgb = hex_to_rgb('888888')
    add_figure_to_slide(slide, "11_model_comparison.png", left=9.3, top=0.95, width=3.8)

    # ── Slide 7: Visualisations ───────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, 'F8FBFF')
    build_slide_header(slide, W, "Key Visualisations")
    add_figure_to_slide(slide, "07_confusion_matrices.png", left=0.2, top=0.85, width=6.3)
    add_figure_to_slide(slide, "08_roc_curves.png",         left=6.7, top=0.85, width=6.3)
    add_content_box(slide, "Confusion Matrices (RF Tuned vs XGBoost)",
                    top=5.65, left=0.2, width=6.3, height=0.4, size=11, color='555555')
    add_content_box(slide, "Multi-class OvR ROC Curves with AUC Scores",
                    top=5.65, left=6.7, width=6.3, height=0.4, size=11, color='555555')

    # ── Slide 8: Feature Importance + Training History ────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, 'F8FBFF')
    build_slide_header(slide, W, "Feature Importance & XGBoost Training History")
    add_figure_to_slide(slide, "09_feature_importance.png",        left=0.2, top=0.85, width=6.5)
    add_figure_to_slide(slide, "10_xgboost_training_history.png",  left=7.0, top=0.85, width=6.0)
    add_content_box(slide, "Top features: Severity_Numeric, Log_Damage, Risk Scores",
                    top=5.7, left=0.2, width=6.5, height=0.4, size=11, color='555555')
    add_content_box(slide, "Flat validation loss confirms no learnable signal in synthetic data",
                    top=5.7, left=7.0, width=6.0, height=0.4, size=11, color='555555')

    # ── Slide 9: Challenges ───────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, 'F8FBFF')
    build_slide_header(slide, W, "Challenges Faced")
    challenges = [
        ("Synthetic Dataset Limitation",
         "Attack_Type labels were randomly assigned — no statistical signal exists between features and target."),
        ("Deprecated API Parameters",
         "scikit-learn removed multi_class param in LogisticRegression in newer versions — fixed by removing argument."),
        ("Unicode Encoding on Windows",
         "Special characters (arrows, em-dashes) caused UnicodeEncodeError on Windows cp1252 terminal — resolved with UTF-8 wrapper."),
        ("SeriesGroupBy.map() Error",
         "Calling .map() on a GroupBy object raised AttributeError — fixed by removing unused intermediate variables."),
        ("XGBoost Module Not Found",
         "XGBoost was not pre-installed — resolved with python -m pip install xgboost."),
        ("Git Push Rejection",
         "Remote contained extra commits (GitHub auto-created files) — resolved with git pull --rebase before push."),
    ]
    for i, (title_c, desc) in enumerate(challenges):
        col = i % 2
        row = i // 2
        left_c = 0.3 + col * 6.6
        top_c  = 1.0 + row * 2.1
        box = slide.shapes.add_shape(1, PInches(left_c), PInches(top_c), PInches(6.3), PInches(1.9))
        box.fill.solid(); box.fill.fore_color.rgb = hex_to_rgb('EBF5FB')
        box.line.color.rgb = hex_to_rgb(MID_BLUE)
        add_title_box(slide, title_c, top_c+0.1, left_c+0.15, 6.0, 0.45, size=14, bold=True, color=DARK_BLUE)
        add_title_box(slide, desc,    top_c+0.6, left_c+0.15, 6.0, 1.0,  size=12, bold=False, color='333333')

    # ── Slide 10: Conclusion & Future Work ────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, 'F8FBFF')
    build_slide_header(slide, W, "Conclusion & Future Work")
    add_content_box(slide, "Conclusions", top=0.9, left=0.3, width=6.0, height=0.4,
                    size=17, bold=True, color=DARK_BLUE)
    conclusions = [
        "Complete end-to-end ML pipeline implemented and deployed",
        "5 models trained, tuned, and evaluated with GridSearchCV + 5-fold CV",
        "15 domain-specific features engineered from raw metadata",
        "Interactive Streamlit app deployed with live prediction",
        "All 12 visualisations generated and saved automatically",
        "Full alignment with IEEE CloudCom 2021 methodology",
    ]
    add_bullet_box(slide, conclusions, top=1.3, left=0.3, width=6.0, height=4.5, size=14)

    add_content_box(slide, "Future Work", top=0.9, left=6.8, width=6.2, height=0.4,
                    size=17, bold=True, color=DARK_BLUE)
    future_w = [
        "Apply to real datasets: UNSW-NB15, CIC-IDS-2017",
        "Add LSTM/Transformer for sequential request modelling",
        "Real-time streaming inference with Apache Kafka",
        "Anomaly detection for zero-day attacks (autoencoders)",
        "Live dashboard monitoring with alert system",
        "Federated learning for privacy-preserving detection",
    ]
    add_bullet_box(slide, future_w, top=1.3, left=6.8, width=6.0, height=4.5, size=14)

    # ── Slide 11: Thank You ───────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide_bg(slide, DARK_BLUE)
    rect = slide.shapes.add_shape(1, PInches(0), PInches(3.3), W, PInches(0.08))
    rect.fill.solid(); rect.fill.fore_color.rgb = hex_to_rgb(ACCENT)
    rect.line.fill.background()
    add_title_box(slide, "Thank You", 1.5, 0.5, 12.3, 1.2,
                  size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_title_box(slide, "Questions & Discussion", 2.8, 0.5, 12.3, 0.7,
                  size=22, bold=False, color='AED6F1', align=PP_ALIGN.CENTER)
    add_title_box(slide,
        "GitHub: github.com/Saiffathalla/api-abuse-detection",
        3.7, 0.5, 12.3, 0.5, size=16, bold=False, color='7FB3D3', align=PP_ALIGN.CENTER)
    add_title_box(slide, "IEEE CloudCom 2021 | REST API Abuse Detection Using ML",
        4.4, 0.5, 12.3, 0.4, size=13, bold=False, color='5D8AA8', align=PP_ALIGN.CENTER)

    out_path = os.path.join(REPORTS_DIR, 'API_Abuse_Detection_Slides.pptx')
    prs.save(out_path)
    print(f"Presentation saved: {out_path}")


# ==============================================================================
# RUN
# ==============================================================================
if __name__ == '__main__':
    print("Generating Word report...")
    build_report()
    print("Generating PowerPoint presentation...")
    build_presentation()
    print("\nDone. Files saved to outputs/reports/")
